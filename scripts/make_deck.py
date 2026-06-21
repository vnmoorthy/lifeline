"""Generate Lifeline_Deck.pptx — a 5-minute pitch deck mapped to the hackathon judging criteria.
    python3 scripts/make_deck.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG = RGBColor(0x0B, 0x0E, 0x14)
CARD = RGBColor(0x14, 0x1C, 0x27)
TEXT = RGBColor(0xEE, 0xF3, 0xF9)
DIM = RGBColor(0x9A, 0xA7, 0xB8)
BLUE = RGBColor(0x5C, 0x9D, 0xFF)
GREEN = RGBColor(0x35, 0xD6, 0xA0)
AMBER = RGBColor(0xF3, 0xB3, 0x40)
RED = RGBColor(0xE5, 0x48, 0x4D)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(tag, title, bullets, accent=BLUE, notes="", big=None, metrics=None):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    # accent bar
    bar = s.shapes.add_shape(1, Inches(0.7), Inches(0.7), Inches(0.18), Inches(0.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    # tag
    tb = s.shapes.add_textbox(Inches(1.0), Inches(0.66), Inches(11), Inches(0.45)).text_frame
    tb.text = tag.upper()
    r = tb.paragraphs[0].runs[0]; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = accent
    # title
    tt = s.shapes.add_textbox(Inches(0.95), Inches(1.15), Inches(11.5), Inches(1.1)).text_frame
    tt.word_wrap = True; tt.text = title
    r = tt.paragraphs[0].runs[0]; r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = TEXT
    y = 2.5
    if big:
        bb = s.shapes.add_textbox(Inches(0.95), Inches(2.4), Inches(11.5), Inches(1.0)).text_frame
        bb.word_wrap = True; bb.text = big
        r = bb.paragraphs[0].runs[0]; r.font.size = Pt(22); r.font.italic = True; r.font.color.rgb = accent
        y = 3.5
    if metrics:
        n = len(metrics); w = 11.6 / n
        for i, (lab, val, col) in enumerate(metrics):
            x = 0.95 + i * w
            box = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w - 0.25), Inches(1.5))
            box.fill.solid(); box.fill.fore_color.rgb = CARD; box.line.color.rgb = RGBColor(0x23, 0x2F, 0x3F); box.line.width = Pt(0.75)
            tf = box.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.18); tf.margin_top = Inches(0.18)
            tf.text = lab
            p0 = tf.paragraphs[0]; p0.runs[0].font.size = Pt(12); p0.runs[0].font.color.rgb = DIM
            p1 = tf.add_paragraph(); p1.text = val; p1.runs[0].font.size = Pt(26); p1.runs[0].font.bold = True; p1.runs[0].font.color.rgb = col
        y += 1.9
    bx = s.shapes.add_textbox(Inches(0.95), Inches(y), Inches(11.5), Inches(7.0 - y)).text_frame
    bx.word_wrap = True
    for i, b in enumerate(bullets):
        p = bx.paragraphs[0] if i == 0 else bx.add_paragraph()
        p.text = "•  " + b; p.space_after = Pt(10)
        p.runs[0].font.size = Pt(18); p.runs[0].font.color.rgb = TEXT
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


# 1 — Title
t = prs.slides.add_slide(BLANK)
t.background.fill.solid(); t.background.fill.fore_color.rgb = BG
mk = t.shapes.add_shape(1, Inches(0.95), Inches(2.3), Inches(0.7), Inches(0.7))
mk.fill.solid(); mk.fill.fore_color.rgb = RED; mk.line.fill.background()
mt = mk.text_frame; mt.text = "+"; mt.paragraphs[0].alignment = PP_ALIGN.CENTER
mt.paragraphs[0].runs[0].font.size = Pt(30); mt.paragraphs[0].runs[0].font.bold = True; mt.paragraphs[0].runs[0].font.color.rgb = TEXT
ti = t.shapes.add_textbox(Inches(1.85), Inches(2.25), Inches(11), Inches(1.0)).text_frame
ti.text = "Lifeline"; ti.paragraphs[0].runs[0].font.size = Pt(54); ti.paragraphs[0].runs[0].font.bold = True; ti.paragraphs[0].runs[0].font.color.rgb = TEXT
sub = t.shapes.add_textbox(Inches(0.95), Inches(3.5), Inches(11.6), Inches(1.4)).text_frame
sub.word_wrap = True
sub.text = "Turn a FROZEN open-source model into a ~98%-verified first-aid assistant — with inference-time compute + a deterministic verifier. Zero training."
sub.paragraphs[0].runs[0].font.size = Pt(22); sub.paragraphs[0].runs[0].font.color.rgb = DIM
tagb = t.shapes.add_textbox(Inches(0.95), Inches(5.4), Inches(11.6), Inches(0.6)).text_frame
tagb.text = "Build the machine  ·  Live app: vnmoorthy.github.io/lifeline  ·  github.com/vnmoorthy/lifeline"
tagb.paragraphs[0].runs[0].font.size = Pt(14); tagb.paragraphs[0].runs[0].font.color.rgb = BLUE
t.notes_slide.notes_text_frame.text = ("Hook: Everyone's instinct for a safety-critical task is fine-tune. We trained nothing. "
    "We took a frozen open model and spent inference-time compute, then verified. Build as if compute is free, then verify.")

# 2 — Vision & Problem Fit
slide("Vision & Problem fit", "In an emergency you can't read a web page — and a wrong instruction can kill",
      ["People Google first-aid in a panic — hands busy, seconds matter, the answer has to be trusted.",
       "Open models are getting good, but 'mostly right' is unacceptable for CPR, choking, overdose, bleeding.",
       "The usual fix — fine-tune on medical data — is slow, costly, and still probabilistic.",
       "Right approach: keep the model frozen, make the SYSTEM safe at inference time."],
      accent=RED,
      notes="Why it matters + why our approach fits: safety-critical, hands-free, and the fix is system design not training.")

# 3 — The thesis (track fit)
slide("The approach", "Don't train — spend. Build as if compute is free, then verify.",
      ["Leave the open model (DiffusionGemma-26B) completely FROZEN — zero gradient steps.",
       "Spend inference-time compute on two knobs: denoising depth × best-of-N.",
       "Gate every candidate through a DETERMINISTIC verifier — keep only protocol-correct answers.",
       "Exactly the 'Build the machine' bet: near-infinite compute, near-zero latency, build for that world today."],
      accent=BLUE,
      big="frozen model  +  inference-time compute  +  a deterministic verifier  =  safe",
      notes="This is the track thesis verbatim. The model stays dumb; the system becomes safe.")

# 4 — How it works
slide("Technical execution", "How it works — recognize → spend compute → verify → speak",
      ["Recognize the spoken emergency (principled triage) → pick the official protocol.",
       "Generate N candidates at the right denoising depth (the effort-manager spends more only when needed).",
       "Verify each against the protocol: required concept-groups present + forbidden actions absent.",
       "Speak the first verified answer — or fall back to the canonical protocol. NEVER an unverified step.",
       "17 emergencies. Hands-free voice. Installable PWA. 50+ tests, all green."],
      accent=BLUE,
      notes="Emphasize what WORKS: tested, runs, the safety invariant (every fallback self-verifies).")

# 5 — The numbers
slide("Technical execution — it works", "Inference-time compute, measured on a frozen model",
      ["Denoising depth lifts a near-random model to its single-sample ceiling (knee at 16 steps).",
       "Best-of-N then closes the gap to ~98% — purely more inference, no training.",
       "Model-agnostic: same recipe on Qwen2.5-7B went 65% → 98.3%."],
      accent=GREEN,
      metrics=[("SINGLE-SHOT", "79.6%", TEXT), ("BEST-OF-N VERIFIED", "98.45%", GREEN),
               ("TRAINING", "None", BLUE), ("LATENCY", "0.5–1.9s", AMBER)],
      notes="The curve: 7.5%→67.5% denoising, 79.6%→98.45% best-of-N. All real, frozen weights.")

# 6 — Novelty & Insight
slide("Novelty & insight", "A deterministic verifier beats an LLM-judge — and that's the whole game",
      ["Best-of-N is adversarial search: with enough samples you WILL surface fluent, confident, wrong answers.",
       "On a 42-candidate adversarial set, the verifier caught 12/12 fluent-but-dangerous answers — 0 unsafe leaks.",
       "An LLM-judge is fooled by fluency, flip-flops run-to-run, and costs an API call per candidate.",
       "Deterministic = consistent, unhackable, ~free (so best-of-N scales), and a hard safety guarantee."],
      accent=AMBER,
      notes="The insight judges remember: more compute only becomes more SAFETY if the selector is a rule, not a vibe.")

# 7 — Demo
slide("Presentation & demo", "Live: talk to it, watch it spend compute, watch it refuse to be wrong",
      ["Easy case ('he collapsed, not breathing') → verifies on try 1, speaks CPR steps instantly.",
       "Hard case ('boiling water on her arm') → escalates best-of-N, early-exits when one verifies.",
       "Money moment: a fluent WRONG answer an LLM-judge would accept → verifier rejects → falls back to protocol.",
       "Live: vnmoorthy.github.io/lifeline (voice app) + vnmoorthy.github.io/lifeline/dashboard.html (results)."],
      accent=BLUE,
      notes="Show the app live (mock_ui), then the dashboard. End on the fluent-wrong rejection.")

# 8 — Impact & Trajectory
slide("Impact & trajectory", "If this works at scale: reliable open models for safety-critical work, no training",
      ["The method is the product: any task with a checkable protocol (code+tests, structured extraction, compliance).",
       "Gets strictly better and cheaper as compute gets cheaper — it rides the curve the hackathon is betting on.",
       "Turns 'a frozen open model' into 'a system you can trust with your life' — without data, labels, or GPUs to train.",
       "Next: more protocols + languages, on-device deployment, live 911 hand-off."],
      accent=GREEN,
      notes="Scale story: it's a general reliability harness for frozen open models, demoed on first aid.")

# 9 — Close
slide("Recap", "Frozen open model + inference-time compute + deterministic verification",
      ["79.6% → 98.45% verified, model-agnostic, zero training.",
       "0 unsafe leaks on the 42-candidate adversarial set — 12/12 fluent-but-dangerous answers rejected.",
       "Live app: vnmoorthy.github.io/lifeline   ·   Dashboard: vnmoorthy.github.io/lifeline/dashboard.html",
       "Code: github.com/vnmoorthy/lifeline   ·   Build as if compute is free, then verify."],
      accent=RED,
      notes="Close confident: this is how you optimize a frozen open model to its limit.")

out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "Lifeline_Deck.pptx")
prs.save(out)
print(f"wrote {out} — {len(prs.slides._sldIdLst)} slides")
